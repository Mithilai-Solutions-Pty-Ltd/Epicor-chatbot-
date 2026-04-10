import os
import time
import json
import hashlib
import logging
import gc
from datetime import datetime, timezone
from typing import List, Dict, Any, Optional

import requests
from dotenv import load_dotenv

load_dotenv()

logger = logging.getLogger("botzi.zoho_sync")


def require_env(name: str) -> str:
    value = os.getenv(name)
    if not value:
        raise RuntimeError(f"Missing required environment variable: {name}")
    return value


def safe_json(resp: requests.Response) -> Dict[str, Any]:
    try:
        return resp.json()
    except Exception:
        return {}


def preview_text(resp: requests.Response, limit: int = 1000) -> str:
    try:
        return resp.text[:limit]
    except Exception:
        return "<unable to read response text>"


def get_zoho_access_token() -> str:
    url = require_env("ZOHO_ACCOUNTS_URL")
    resp = requests.post(
        url,
        data={
            "grant_type": "refresh_token",
            "client_id": require_env("ZOHO_CLIENT_ID"),
            "client_secret": require_env("ZOHO_CLIENT_SECRET"),
            "refresh_token": require_env("ZOHO_REFRESH_TOKEN"),
        },
        timeout=30,
    )
    resp.raise_for_status()
    token = safe_json(resp).get("access_token")
    if not token:
        raise ValueError(f"No access_token in Zoho response: {preview_text(resp)}")
    logger.info("✅ Zoho access token obtained")
    return token


def build_headers(access_token: str, accept: str = "application/vnd.api+json") -> Dict[str, str]:
    return {
        "Authorization": f"Zoho-oauthtoken {access_token}",
        "Accept": accept,
    }


def list_folder_recursive(folder_id: str, access_token: str, depth: int = 0) -> List[Dict[str, Any]]:
    base_url = os.environ["ZOHO_WORKDRIVE_URL"].rstrip("/")
    headers = {
        "Authorization": f"Zoho-oauthtoken {access_token}",
        "Accept": "application/vnd.api+json",
    }
    files: List[Dict[str, Any]] = []
    if depth == 0:
        endpoint = f"{base_url}/teamfolders/{folder_id}/files"
    else:
        endpoint = f"{base_url}/files/{folder_id}/files"

    resp = requests.get(
        endpoint,
        headers=headers,
        params={"page[limit]": 200, "page[offset]": 0},
        timeout=30,
    )
    if not resp.ok:
        logger.error(f"WorkDrive API error {resp.status_code}: {resp.text[:500]}")
    resp.raise_for_status()

    items = resp.json().get("data", [])
    logger.info(f"{'  ' * depth}Found {len(items)} items in folder {folder_id}")

    for item in items:
        attrs = item.get("attributes", {}) or {}
        top_type = str(item.get("type", "") or "").strip().lower()
        item_type = str(attrs.get("type", "") or "").strip().lower()
        resource = str(attrs.get("resource_type", "") or "").strip().lower()
        file_name = str(attrs.get("name", "") or "")
        file_id = str(item.get("id", "") or "")
        modified = str(attrs.get("modified_time", "") or "")

        is_folder = (
            top_type == "folders"
            or item_type == "folder"
            or resource == "folder"
        )

        if is_folder:
            logger.info(f"{'  ' * depth}📁 Scanning folder: {file_name}")
            sub_files = list_folder_recursive(file_id, access_token, depth + 1)
            files.extend(sub_files)
            continue

        ext = file_name.lower().rsplit(".", 1)[-1] if "." in file_name else ""
        if ext in ("pdf", "pptx", "docx", "txt", "mp4"):
            logger.info(f"{'  ' * depth}📄 Found: {file_name} ({ext})")
            files.append({
                "file_id": file_id,
                "file_name": file_name,
                "doc_type": ext,
                "modified": modified,
                "folder_depth": depth,
            })
        else:
            logger.debug(f"{'  ' * depth}⏭️ Skipping unsupported: {file_name}")

    return files


def extract_download_url_from_metadata(meta: Dict[str, Any]) -> Optional[str]:
    data = meta.get("data", {}) or {}
    attrs = data.get("attributes", {}) or {}
    links = data.get("links", {}) or {}
    top_links = meta.get("links", {}) or {}
    candidates = [
        attrs.get("download_url"),
        attrs.get("downloadUrl"),
        attrs.get("Download URL"),
        links.get("download"),
        links.get("self"),
        top_links.get("download"),
    ]
    for candidate in candidates:
        if isinstance(candidate, str) and candidate.startswith("http"):
            return candidate
    return None


def try_binary_download(url: str, access_token: str) -> Optional[bytes]:
    try:
        resp = requests.get(
            url,
            headers={
                "Authorization": f"Zoho-oauthtoken {access_token}",
                "Accept": "*/*",
            },
            timeout=120,
            allow_redirects=True,
            stream=False,
        )
    except Exception as e:
        logger.warning("Download request failed for %s: %s", url, e)
        return None

    content_type = (resp.headers.get("Content-Type") or "").lower()
    logger.info("    GET %s -> %s (%s)", url, resp.status_code, content_type)

    if resp.ok and len(resp.content) > 100 and "application/json" not in content_type:
        logger.info("    Downloaded %d bytes", len(resp.content))
        return resp.content

    if "application/json" in content_type:
        body = safe_json(resp)
        logger.warning("    JSON response from download URL: %s", json.dumps(body)[:1000])
        nested_url = extract_download_url_from_metadata(body)
        if nested_url and nested_url != url:
            logger.info("    Following nested download URL")
            return try_binary_download(nested_url, access_token)

    if not resp.ok:
        logger.warning("    Download failed: status=%s body=%s", resp.status_code, preview_text(resp))

    return None


def get_preview_download_url(file_id: str, access_token: str) -> Optional[str]:
    base_url = require_env("ZOHO_WORKDRIVE_URL").rstrip("/")
    headers = build_headers(access_token)
    resp = requests.get(f"{base_url}/files/{file_id}/previewinfo", headers=headers, timeout=30)
    if not resp.ok:
        logger.warning("    previewinfo failed: %s %s", resp.status_code, preview_text(resp))
        return None
    url = safe_json(resp).get("data", {}).get("attributes", {}).get("preview_data_url")
    if url:
        logger.info("    Got preview_data_url: %s", url)
    return url


def download_file(file_id: str, access_token: str) -> bytes:
    base_url = require_env("ZOHO_WORKDRIVE_URL").rstrip("/")
    headers = build_headers(access_token)

    info_url = f"{base_url}/files/{file_id}"
    info_resp = requests.get(info_url, headers=headers, timeout=30)
    logger.info("    Fetching metadata for file_id=%s -> %s", file_id, info_resp.status_code)

    if not info_resp.ok:
        raise ValueError(
            f"Metadata fetch failed for file_id={file_id}, "
            f"status={info_resp.status_code}, body={preview_text(info_resp)}"
        )

    info_json = safe_json(info_resp)
    attrs = info_json.get("data", {}).get("attributes", {}) or {}
    logger.info("    Metadata keys: %s", list(attrs.keys()))

    candidate_urls = []
    metadata_download_url = extract_download_url_from_metadata(info_json)
    if metadata_download_url:
        candidate_urls.append(metadata_download_url)
    candidate_urls.append(f"{base_url}/files/{file_id}/download")

    seen = set()
    deduped_urls = []
    for url in candidate_urls:
        if url and url not in seen:
            deduped_urls.append(url)
            seen.add(url)

    for url in deduped_urls:
        logger.info("    Trying download URL: %s", url)
        content = try_binary_download(url, access_token)
        if content:
            return content

    logger.warning("    Direct download failed, falling back to preview engine")
    preview_url = get_preview_download_url(file_id, access_token)
    if preview_url:
        content = try_binary_download(preview_url, access_token)
        if content:
            logger.info("    Downloaded via preview engine (%d bytes)", len(content))
            return content

    raise ValueError(
        f"Download failed for file_id={file_id}. "
        f"Tried: {deduped_urls} + preview engine"
    )


def extract_text_from_pdf(data: bytes) -> List[Dict[str, Any]]:
    import fitz
    pages = []
    try:
        doc = fitz.open(stream=data, filetype="pdf")
        for i in range(len(doc)):
            page = doc[i]
            text = page.get_text("text").strip()
            if text:
                pages.append({"text": text, "page": i + 1})
            page = None
            if i % 50 == 0:
                gc.collect()
        doc.close()
        del doc
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
    finally:
        gc.collect()
    return pages


def extract_text_from_pptx(data: bytes) -> List[Dict[str, Any]]:
    from io import BytesIO
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append({"text": "\n".join(texts), "page": i})
    return slides


def extract_text_from_docx(data: bytes) -> List[Dict[str, Any]]:
    from io import BytesIO
    from docx import Document
    doc = Document(BytesIO(data))
    chunks_out = []
    page, buffer, word_count = 1, [], 0
    for para in doc.paragraphs:
        txt = para.text.strip()
        if not txt:
            continue
        buffer.append(txt)
        word_count += len(txt.split())
        if word_count >= 500:
            chunks_out.append({"text": "\n".join(buffer), "page": page})
            page += 1
            buffer, word_count = [], 0
    if buffer:
        chunks_out.append({"text": "\n".join(buffer), "page": page})
    return chunks_out


def extract_text_from_txt(data: bytes) -> List[Dict[str, Any]]:
    text = data.decode("utf-8", errors="ignore")
    lines = text.split("\n")
    chunks_out, buffer, page = [], [], 1
    for line in lines:
        buffer.append(line)
        if len(buffer) >= 80:
            chunks_out.append({"text": "\n".join(buffer), "page": page})
            page += 1
            buffer = []
    if buffer:
        chunks_out.append({"text": "\n".join(buffer), "page": page})
    return chunks_out


def extract_text(data: bytes, doc_type: str) -> List[Dict[str, Any]]:
    extractor_map = {
        "pdf":  extract_text_from_pdf,
        "pptx": extract_text_from_pptx,
        "docx": extract_text_from_docx,
        "txt":  extract_text_from_txt,
        "mp4":  lambda d: [],
    }
    fn = extractor_map.get(doc_type)
    if fn is None:
        return []
    return fn(data)


def split_into_chunks(
    pages: List[Dict[str, Any]],
    file_id: str,
    file_name: str,
    doc_type: str,
    chunk_size: int = 800,
    overlap: int = 100,
) -> List[Dict[str, Any]]:
    chunks = []
    for page_data in pages:
        text = page_data["text"]
        page = page_data["page"]
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            snippet = text[start:end].strip()
            if len(snippet) > 50:
                chunk_id = hashlib.md5(f"{file_id}|{page}|{start}".encode()).hexdigest()
                chunks.append({
                    "id": chunk_id,
                    "text": snippet,
                    "metadata": {
                        "file_id":   file_id,
                        "file_name": file_name,
                        "doc_type":  doc_type,
                        "source":    file_name,
                        "page":      int(page),
                    },
                })
            start += chunk_size - overlap
    return chunks


def get_sync_log() -> Dict[str, str]:
    from supabase import create_client
    sb = create_client(require_env("SUPABASE_URL"), require_env("SUPABASE_KEY"))
    result = sb.table("sync_log").select("file_id,modified").execute()
    return {row["file_id"]: row["modified"] for row in (result.data or [])}


def update_sync_log(file_id: str, file_name: str, modified: str, chunks: int):
    from supabase import create_client
    sb = create_client(require_env("SUPABASE_URL"), require_env("SUPABASE_KEY"))
    sb.table("sync_log").upsert({
        "file_id":   file_id,
        "file_name": file_name,
        "modified":  modified,
        "chunks":    chunks,
        "synced_at": datetime.now(timezone.utc).isoformat(),
    }).execute()


# ── NEW: Stream download directly to disk ─────────────────────────────────────
def download_file_to_path(file_id: str, access_token: str, dest_path: str) -> None:
    """Stream a WorkDrive file directly to disk — never holds full bytes in RAM."""
    base_url = require_env("ZOHO_WORKDRIVE_URL").rstrip("/")
    headers = build_headers(access_token)

    info_resp = requests.get(f"{base_url}/files/{file_id}", headers=headers, timeout=30)
    info_resp.raise_for_status()
    info_json = safe_json(info_resp)

    candidate_urls = []
    metadata_url = extract_download_url_from_metadata(info_json)
    if metadata_url:
        candidate_urls.append(metadata_url)
    candidate_urls.append(f"{base_url}/files/{file_id}/download")

    for url in candidate_urls:
        logger.info("    Streaming download: %s", url)
        try:
            resp = requests.get(
                url,
                headers={**headers, "Accept": "*/*"},
                timeout=120,
                allow_redirects=True,
                stream=True,                        # ← key: stream, don't buffer
            )
        except Exception as e:
            logger.warning("    Stream request failed: %s", e)
            continue

        content_type = (resp.headers.get("Content-Type") or "").lower()

        if resp.ok and "application/json" not in content_type:
            total = 0
            with open(dest_path, "wb") as f:
                for chunk in resp.iter_content(chunk_size=65536):  # 64KB at a time
                    if chunk:
                        f.write(chunk)
                        total += len(chunk)
            logger.info("    Written %d bytes to disk", total)
            if total > 100:
                return  # success

        elif "application/json" in content_type:
            body = safe_json(resp)
            nested_url = extract_download_url_from_metadata(body)
            if nested_url and nested_url not in candidate_urls:
                candidate_urls.append(nested_url)

    raise ValueError(f"Stream download failed for file_id={file_id}")


# ── NEW: Extract PDF text from disk path (low RAM) ────────────────────────────
def _extract_pdf_from_path(pdf_path: str) -> List[Dict[str, Any]]:
    """Open PDF from disk — fitz reads pages lazily, avoids bytes+doc in RAM together."""
    import fitz
    pages = []
    try:
        doc = fitz.open(pdf_path)                  # ← path, not stream=bytes
        for i in range(len(doc)):
            page = doc[i]
            text = page.get_text("text").strip()
            if text:
                pages.append({"text": text, "page": i + 1})
            page = None
            if i % 50 == 0:
                gc.collect()
        doc.close()
        del doc
    except Exception as e:
        logger.error("PDF extraction error: %s", e)
    finally:
        gc.collect()
    return pages


# ── NEW: Unified path-based extractor ────────────────────────────────────────
def extract_text_from_path(file_path: str, doc_type: str) -> List[Dict[str, Any]]:
    """Extract text from a file already on disk."""
    if doc_type == "pdf":
        return _extract_pdf_from_path(file_path)
    if doc_type in ("pptx", "docx", "txt"):
        with open(file_path, "rb") as f:
            data = f.read()
        result = extract_text(data, doc_type)
        del data
        gc.collect()
        return result
    return []


def run_sync():
    import subprocess
    import sys

    logger.info("=" * 60)
    logger.info("🔄 BOTZI Zoho WorkDrive Sync Started")
    logger.info("=" * 60)

    token = get_zoho_access_token()
    team_folder_id = require_env("ZOHO_TEAM_FOLDER_ID")
    logger.info("Scanning team folder: %s", team_folder_id)

    all_files = list_folder_recursive(team_folder_id, token)
    logger.info("Found %d supported files in WorkDrive", len(all_files))

    sync_log = get_sync_log()
    logger.info("Previously synced: %d files", len(sync_log))

    new_count     = 0
    updated_count = 0
    skipped_count = 0

    for file_info in all_files:
        file_id   = file_info["file_id"]
        file_name = file_info["file_name"]
        modified  = file_info["modified"]
        doc_type  = file_info["doc_type"]

        # ── Skip MP4 always — no text ──────────────────────────
        if doc_type == "mp4":
            skipped_count += 1
            logger.info("  ⏭️  Skipping mp4: %s", file_name)
            continue

        # ── Skip unchanged files ───────────────────────────────
        if file_id in sync_log:
            if sync_log[file_id] == modified:
                skipped_count += 1
                logger.info("  ⏭️  Skipping unchanged: %s", file_name)
                continue
            else:
                logger.info("  🔄 Re-indexing modified: %s", file_name)
                is_update = True
        else:
            logger.info("  🆕 New file: %s", file_name)
            is_update = False

        # ── Spawn fresh subprocess per file ────────────────────
        try:
            result = subprocess.run(
                [
                    sys.executable,
                    "-m", "zoho_sync.sync_one_file",
                    file_id,
                    file_name,
                    modified,
                    doc_type,
                    "true" if is_update else "false",
                ],
                timeout=600,
                capture_output=False,
            )

            if result.returncode == 0:
                if is_update:
                    updated_count += 1
                else:
                    new_count += 1
                logger.info("  ✅ Completed: %s", file_name)
            else:
                logger.error(
                    "  ❌ Failed (exit %d): %s",
                    result.returncode, file_name
                )

        except subprocess.TimeoutExpired:
            logger.error("  ❌ Timeout (10min): %s", file_name)
        except Exception as e:
            logger.error("  ❌ Error: %s — %s", file_name, e)

        time.sleep(2)

    logger.info("=" * 60)
    logger.info(
        "✅ Sync complete: %d new | %d updated | %d skipped",
        new_count, updated_count, skipped_count,
    )
    logger.info("=" * 60)


if __name__ == "__main__":
    run_sync()